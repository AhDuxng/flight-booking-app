import sys
import os
import subprocess
# Tự động cài đặt python-pptx nếu chưa có
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Thư viện 'python-pptx' chưa được cài đặt. Đang cài đặt tự động...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
def create_presentation():
    prs = Presentation()
    
    # Thiết lập kích thước slide 16:9 tiêu chuẩn cho Canva
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Định nghĩa màu sắc thương hiệu Phenikaa
    COLOR_BLUE = RGBColor(24, 43, 97)       # Xanh đậm Phenikaa
    COLOR_ORANGE = RGBColor(243, 112, 33)   # Cam Phenikaa
    COLOR_WHITE = RGBColor(255, 255, 255)   # Trắng
    COLOR_GRAY = RGBColor(241, 245, 249)    # Xám nền nhẹ
    COLOR_TEXT_DARK = RGBColor(15, 23, 42)  # Chữ tối
    COLOR_ALERT = RGBColor(220, 38, 38)     # Màu đỏ cảnh báo
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Helper: Tạo slide tiêu đề hoặc kết thúc (Nền xanh đậm)
    def add_title_slide(title_text, subtitle_text, author_text):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Nền xanh đậm
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BLUE
        bg.line.fill.background()
        
        # Tiêu đề chính
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(44)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Tiêu đề phụ
        txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.333), Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(24)
        p2.font.color.rgb = COLOR_ORANGE
        p2.alignment = PP_ALIGN.CENTER
        
        # Thông tin tác giả
        txBox3 = slide.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.333), Inches(1.5))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = author_text
        p3.font.size = Pt(16)
        p3.font.color.rgb = COLOR_WHITE
        p3.alignment = PP_ALIGN.CENTER
        
        return slide
    # Helper: Tạo slide nội dung (Nền xám nhẹ, có Header xanh đậm, Accent cam)
    def add_content_slide(title_text):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Nền xám nhẹ
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_GRAY
        bg.line.fill.background()
        
        # Header xanh đậm
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_BLUE
        header.line.fill.background()
        
        # Chữ tiêu đề
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.2), Inches(11.833), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.LEFT
        
        # Dải màu cam trang trí dưới header
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), Inches(13.333), Inches(0.1)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = COLOR_ORANGE
        stripe.line.fill.background()
        
        return slide
    # ==========================================
    # SLIDE 1: Trang tiêu đề
    # ==========================================
    add_title_slide(
        "HỆ THỐNG ĐẶT VÉ MÁY BAY TRỰC TUYẾN",
        "Báo cáo Kiến trúc & Giải pháp Công nghệ",
        "Nhóm 6\nPhạm Anh Dũng   -   Trịnh Đức Anh   -   Phạm Thiên Ân\n\nTrường Đại học Phenikaa"
    )
    # ==========================================
    # SLIDE 2: Mục lục
    # ==========================================
    slide2 = add_content_slide("NỘI DUNG BÁO CÁO")
    col1 = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.0), Inches(4.5))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    
    bullets1 = [
        "1. Giới thiệu & Đặt vấn đề",
        "2. Kiến trúc & Công nghệ Sử dụng",
        "3. Thiết kế Cơ sở Dữ liệu (Database)",
    ]
    for i, bullet in enumerate(bullets1):
        p = tf1.add_paragraph() if i > 0 else tf1.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(20)
        
    col2 = slide2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.0), Inches(4.5))
    tf2 = col2.text_frame
    tf2.word_wrap = True
    
    bullets2 = [
        "4. Các giải pháp kỹ thuật nâng cao",
        "5. Phân hệ tính năng Client & Admin",
        "6. Hiện trạng & Kế hoạch phát triển",
    ]
    for i, bullet in enumerate(bullets2):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = bullet
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(20)
    # ==========================================
    # SLIDE 3: Bối cảnh & Thách thức
    # ==========================================
    slide3 = add_content_slide("Bối cảnh & Thách thức nghiệp vụ")
    box1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
    box1.fill.solid()
    box1.fill.fore_color.rgb = COLOR_WHITE
    box1.line.color.rgb = COLOR_BLUE
    box1.line.width = Pt(1.5)
    
    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thách thức kỹ thuật trong bài toán đặt vé máy bay:"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_BLUE
    p.space_after = Pt(15)
    
    bullets = [
        "• Sự tăng trưởng nhu cầu: Yêu cầu hệ thống phản hồi tức thời, tra cứu mượt mà.",
        "• Race Condition (Tranh chấp ghế): Tránh trường hợp nhiều khách hàng thanh toán cùng một ghế cùng lúc.",
        "• Giữ ghế ảo đầu cơ: Khách hàng click chọn giữ ghế quá lâu nhưng không thực hiện thanh toán.",
        "• Tính nhất quán dữ liệu & tài chính: Giá vé thay đổi liên tục, cần snapshot lại giá tại thời điểm đặt."
    ]
    for bullet in bullets:
        p2 = tf.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_after = Pt(10)
    # ==========================================
    # SLIDE 4: Mục tiêu & Giải pháp tổng thể
    # ==========================================
    slide4 = add_content_slide("Mục tiêu & Giải pháp tổng thể")
    solutions = [
        ("1. Kiến trúc Decoupled", "Tách độc lập Frontend (React) và Backend API (Express), dễ bảo trì và phân tách trách nhiệm nghiệp vụ rõ ràng.", Inches(0.8)),
        ("2. Trải nghiệm đặt vé mượt", "Quy trình đặt vé tối giản 6 bước khép kín, trực quan hóa sơ đồ chọn ghế trên tàu bay.", Inches(4.8)),
        ("3. Bảo vệ từ Database", "Giải quyết tranh chấp chỗ ngồi, mã giảm giá và bảo mật RLS trực tiếp tại tầng cơ sở dữ liệu.", Inches(8.8)),
    ]
    
    for title, desc, left in solutions:
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(3.7), Inches(4.0))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_WHITE
        box.line.color.rgb = COLOR_ORANGE
        box.line.width = Pt(2)
        
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_BLUE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(15)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.alignment = PP_ALIGN.LEFT
    # ==========================================
    # SLIDE 5: Kiến trúc hệ thống
    # ==========================================
    slide5 = add_content_slide("Kiến trúc hệ thống API-First")
    
    client_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.2), Inches(3.0), Inches(1.8))
    client_box.fill.solid()
    client_box.fill.fore_color.rgb = COLOR_WHITE
    client_box.line.color.rgb = COLOR_BLUE
    client_box.line.width = Pt(2)
    tf = client_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "React Client\n(Frontend)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    api_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(2.2), Inches(3.0), Inches(1.8))
    api_box.fill.solid()
    api_box.fill.fore_color.rgb = COLOR_WHITE
    api_box.line.color.rgb = COLOR_BLUE
    api_box.line.width = Pt(2)
    tf = api_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Node.js Express\n(Backend API)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    db_box = slide5.shapes.add_shape(MSO_SHAPE.CAN, Inches(9.3), Inches(2.2), Inches(3.0), Inches(1.8))
    db_box.fill.solid()
    db_box.fill.fore_color.rgb = COLOR_WHITE
    db_box.line.color.rgb = COLOR_ORANGE
    db_box.line.width = Pt(2)
    tf = db_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "PostgreSQL\n(Supabase DB)"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    arrow1 = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.1), Inches(2.8), Inches(1.0), Inches(0.4))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLOR_BLUE
    arrow1.line.fill.background()
    
    arrow2 = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.3), Inches(2.8), Inches(0.9), Inches(0.4))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLOR_ORANGE
    arrow2.line.fill.background()
    
    desc_box = slide5.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.333), Inches(2.0))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Mô tả luồng hoạt động:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_BLUE
    p.space_after = Pt(10)
    
    bullets = [
        "1. Client thực hiện gửi request API kèm JWT Token qua Axios.",
        "2. Backend tiếp nhận, chạy middleware kiểm tra bảo mật (Helmet, Rate Limiter) và xác thực JWT.",
        "3. Backend truy vấn an toàn và thực hiện transaction xử lý nghiệp vụ với PostgreSQL."
    ]
    for bullet in bullets:
        p2 = tf.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_TEXT_DARK
    # ==========================================
    # SLIDE 6: Công nghệ Sử dụng
    # ==========================================
    slide6 = add_content_slide("Công nghệ cốt lõi trong hệ thống")
    table_shape = slide6.shapes.add_table(6, 2, Inches(1.5), Inches(2.0), Inches(10.333), Inches(4.5))
    table = table_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(7.333)
    
    headers = ["Phân hệ / Lớp", "Công nghệ & Thư viện sử dụng"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = COLOR_WHITE
            
    rows = [
        ["Frontend (Client)", "React 18, Vite, Tailwind CSS v4, shadcn/ui, Lucide Icons"],
        ["Backend (API Server)", "Node.js, Express, Zod (validation), Nodemailer (gửi hóa đơn)"],
        ["Database & Auth", "Supabase (PostgreSQL), pg_cron (tác vụ ngầm)"],
        ["Thanh toán tích hợp", "VNPay Sandbox, MoMo, Stripe Sandbox (Simulated)"],
        ["Bảo mật hệ thống", "JWT, bcryptjs, Helmet, Express-Rate-Limit"],
    ]
    
    for row_idx, row_data in enumerate(rows):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = COLOR_TEXT_DARK
    # ==========================================
    # SLIDE 7: Thiết kế Cơ sở dữ liệu
    # ==========================================
    slide7 = add_content_slide("Cơ sở dữ liệu - Nguyên tắc chuẩn hóa")
    box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BLUE
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Cơ sở dữ liệu quan hệ PostgreSQL trên nền tảng Supabase:"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_BLUE
    p.space_after = Pt(15)
    
    bullets = [
        "• Toàn vẹn dữ liệu (ACID): Đảm bảo các giao dịch đặt chỗ và thanh toán diễn ra nhất quán.",
        "• Bảo mật thông tin: Sử dụng khóa chính ngẫu nhiên UUID (Universal Unique Identifier) thay thế ID số tăng dần, ngăn chặn tấn công dò quét thông tin.",
        "• An toàn tài chính: Cột giá tiền sử dụng kiểu dữ liệu NUMERIC(12,2) để triệt tiêu lỗi làm tròn dấu phẩy động.",
        "• Chỉ mục tối ưu (Database Indexes): Tạo Index trên các trường tìm kiếm chính như departure_time, origin, destination, status để đẩy nhanh tốc độ truy vấn."
    ]
    for bullet in bullets:
        p2 = tf.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_after = Pt(10)
    # ==========================================
    # SLIDE 8: 19 Bảng Dữ liệu
    # ==========================================
    slide8 = add_content_slide("Cấu trúc cơ sở dữ liệu quan hệ (19 Bảng)")
    box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_ORANGE
    box.line.width = Pt(1.5)
    
    title_box = slide8.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.333), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = "Danh sách phân hệ bảng dữ liệu được chuẩn hóa:"
    p_title.font.bold = True
    p_title.font.size = Pt(20)
    p_title.font.color.rgb = COLOR_BLUE
    
    col1 = slide8.shapes.add_textbox(Inches(1.5), Inches(2.7), Inches(5.0), Inches(4.0))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    
    items1 = [
        "1. Vận hành đường bay:",
        "   - airlines (Hãng bay)",
        "   - airports (Sân bay)",
        "   - aircrafts (Tàu bay)",
        "   - flights (Chuyến bay)",
        "2. Nghiệp vụ đặt vé:",
        "   - users (Thông tin mở rộng)",
        "   - bookings (Đơn đặt vé)",
        "   - passengers (Hành khách)"
    ]
    for i, item in enumerate(items1):
        p = tf1.add_paragraph() if i > 0 else tf1.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        if ":" in item:
            p.font.bold = True
            p.font.color.rgb = COLOR_ORANGE
        else:
            p.font.color.rgb = COLOR_TEXT_DARK
            
    col2 = slide8.shapes.add_textbox(Inches(7.0), Inches(2.7), Inches(5.0), Inches(4.0))
    tf2 = col2.text_frame
    tf2.word_wrap = True
    
    items2 = [
        "3. Dịch vụ phụ trợ & Ưu đãi:",
        "   - baggage_options / booking_baggage",
        "   - meal_options / booking_meals",
        "   - discounts / booking_discounts",
        "4. Quản trị & Hệ thống:",
        "   - payments (Nhật ký thanh toán)",
        "   - notifications (Thông báo realtime)",
        "   - reviews (Phản hồi người dùng)",
        "   - admin_logs (Lưu vết hành động)"
    ]
    for i, item in enumerate(items2):
        p = tf2.add_paragraph() if i > 0 else tf2.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        if ":" in item:
            p.font.bold = True
            p.font.color.rgb = COLOR_ORANGE
        else:
            p.font.color.rgb = COLOR_TEXT_DARK
    # ==========================================
    # SLIDE 9: Chống Race Condition đặt ghế
    # ==========================================
    slide9 = add_content_slide("Giải pháp 1: Chống Race Condition đặt ghế")
    user1 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.2), Inches(2.8), Inches(1.2))
    user1.fill.solid()
    user1.fill.fore_color.rgb = COLOR_WHITE
    user1.line.color.rgb = COLOR_BLUE
    user1.line.width = Pt(1.5)
    user1.text = "Khách hàng A\nĐặt ghế 12A"
    user1.text_frame.paragraphs[0].font.color.rgb = COLOR_BLUE
    user1.text_frame.paragraphs[0].font.size = Pt(16)
    
    user2 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.8), Inches(2.8), Inches(1.2))
    user2.fill.solid()
    user2.fill.fore_color.rgb = COLOR_WHITE
    user2.line.color.rgb = COLOR_BLUE
    user2.line.width = Pt(1.5)
    user2.text = "Khách hàng B\nĐặt ghế 12A"
    user2.text_frame.paragraphs[0].font.color.rgb = COLOR_BLUE
    user2.text_frame.paragraphs[0].font.size = Pt(16)
    
    lock_shape = slide9.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(4.8), Inches(2.5), Inches(3.2), Inches(2.2))
    lock_shape.fill.solid()
    lock_shape.fill.fore_color.rgb = COLOR_ORANGE
    lock_shape.line.fill.background()
    
    tf = lock_shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    
    p1 = tf.paragraphs[0]
    p1.text = "PostgreSQL Lock"
    p1.font.bold = True
    p1.font.size = Pt(13)
    p1.font.color.rgb = COLOR_WHITE
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "SELECT FOR UPDATE"
    p2.font.bold = True
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER
    
    res1 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(2.2), Inches(3.2), Inches(1.2))
    res1.fill.solid()
    res1.fill.fore_color.rgb = COLOR_WHITE
    res1.line.color.rgb = COLOR_BLUE
    res1.line.width = Pt(1.5)
    res1.text = "Giao dịch A:\nGiữ ghế thành công ('held')"
    res1.text_frame.paragraphs[0].font.color.rgb = COLOR_BLUE
    res1.text_frame.paragraphs[0].font.size = Pt(14)
    
    res2 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(3.8), Inches(3.2), Inches(1.2))
    res2.fill.solid()
    res2.fill.fore_color.rgb = COLOR_WHITE
    res2.line.color.rgb = COLOR_ALERT
    res2.line.width = Pt(1.5)
    res2.text = "Giao dịch B:\nBị từ chối (Seat not available)"
    res2.text_frame.paragraphs[0].font.color.rgb = COLOR_ALERT
    res2.text_frame.paragraphs[0].font.size = Pt(14)
    
    arr1 = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.0), Inches(2.7), Inches(0.7), Inches(0.2))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = COLOR_BLUE
    arr1.line.fill.background()
    
    arr2 = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.0), Inches(4.3), Inches(0.7), Inches(0.2))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = COLOR_BLUE
    arr2.line.fill.background()
    
    arr3 = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.2), Inches(2.7), Inches(0.8), Inches(0.2))
    arr3.fill.solid()
    arr3.fill.fore_color.rgb = COLOR_ORANGE
    arr3.line.fill.background()
    
    arr4 = slide9.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.2), Inches(4.3), Inches(0.8), Inches(0.2))
    arr4.fill.solid()
    arr4.fill.fore_color.rgb = COLOR_ORANGE
    arr4.line.fill.background()
    
    desc = slide9.shapes.add_textbox(Inches(1.0), Inches(5.4), Inches(11.333), Inches(1.5))
    tf = desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• Giao dịch đầu tiên thực hiện hold_seat() sẽ khóa dòng dữ liệu ghế lại."
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT_DARK
    p2 = tf.add_paragraph()
    p2.text = "• Giao dịch thứ hai cố gắng truy cập sẽ bị chặn lại và nhận mã lỗi từ Database, chống đặt trùng."
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_TEXT_DARK
    # ==========================================
    # SLIDE 10: Tự động giải phóng ghế held
    # ==========================================
    slide10 = add_content_slide("Giải pháp 2: Giải phóng ghế ảo bằng pg_cron")
    box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BLUE
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Quy trình giải phóng tài nguyên ghế tự động:"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_BLUE
    p.space_after = Pt(15)
    
    bullets = [
        "1. Khóa tạm thời: Khi chọn ghế, ghế chuyển sang trạng thái held trong vòng 10 phút để người dùng thanh toán.",
        "2. Tác vụ quét ngầm: Sử dụng Postgres function release_expired_held_seats() kiểm tra thời gian cập nhật.",
        "3. Tự động hóa qua pg_cron: Lập lịch cho hệ thống quét định kỳ mỗi 5 phút một lần để giải phóng các ghế giữ quá hạn về lại available.",
        "4. Hiệu quả: Tối đa hóa công suất ghế thực tế, loại bỏ triệt để hành vi đầu cơ hoặc giữ ghế ảo."
    ]
    for bullet in bullets:
        p2 = tf.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_after = Pt(10)
    # ==========================================
    # SLIDE 11: Bảo mật RLS & Realtime
    # ==========================================
    slide11 = add_content_slide("Giải pháp 3: Bảo mật RLS & Đồng bộ Realtime")
    box = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BLUE
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Phân quyền nâng cao (Row Level Security) & Realtime:"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_BLUE
    p.space_after = Pt(15)
    
    bullets = [
        "• Row Level Security (RLS): Bật RLS cho toàn bộ các bảng. Khách hàng chỉ có quyền xem/sửa các bản ghi của chính họ (user_id = auth.uid()). Tránh rò rỉ dữ liệu người dùng.",
        "• Phân quyền Admin: Quyền admin được tự động xác minh ở tầng database bằng cách đọc metadata (role = 'admin') trong chuỗi mã hóa JWT token.",
        "• Đồng bộ Realtime: Sử dụng cơ chế lắng nghe sự kiện của Supabase Realtime cho bảng bookings và notifications giúp cập nhật trạng thái giao dịch và gửi thông báo tức thời cho người dùng."
    ]
    for bullet in bullets:
        p2 = tf.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_after = Pt(10)
    # ==========================================
    # SLIDE 12: Luồng tính năng người dùng
    # ==========================================
    slide12 = add_content_slide("Quy trình 6 bước phía Khách hàng (User)")
    steps = [
        ("1. Tìm", Inches(0.8)),
        ("2. Điền", Inches(2.8)),
        ("3. Ghế", Inches(4.8)),
        ("4. Dịch vụ", Inches(6.8)),
        ("5. Áp mã", Inches(8.8)),
        ("6. Trả", Inches(10.8)),
    ]
    
    for title, left in steps:
        step_circle = slide12.shapes.add_shape(MSO_SHAPE.OVAL, left, Inches(2.2), Inches(1.6), Inches(1.6))
        step_circle.fill.solid()
        step_circle.fill.fore_color.rgb = COLOR_BLUE
        step_circle.line.color.rgb = COLOR_ORANGE
        step_circle.line.width = Pt(2)
        tf = step_circle.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        if left != Inches(10.8):
            arr = slide12.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(1.6), Inches(2.8), Inches(0.4), Inches(0.2))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR_ORANGE
            arr.line.fill.background()
            
    desc = slide12.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.5))
    tf = desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Chi tiết các chức năng người dùng:"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_BLUE
    p.space_after = Pt(10)
    
    bullets = [
        "• Đầy đủ các bước chọn điểm khởi hành/đến, lọc giá vé, xem chi tiết đường bay.",
        "• Sơ đồ ghế ngồi trực quan hiển thị trạng thái ghế trống (economy, business, first class).",
        "• Mua kèm hành lý ký gửi bổ sung, chọn các suất ăn tiêu chuẩn và nhập mã coupon giảm giá."
    ]
    for bullet in bullets:
        p2 = tf.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_TEXT_DARK
    # ==========================================
    # SLIDE 13: Tính năng của Admin
    # ==========================================
    slide13 = add_content_slide("Hệ thống quản trị (Admin Dashboard)")
    box = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_WHITE
    box.line.color.rgb = COLOR_BLUE
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Các tính năng nghiệp vụ của quản trị viên:"
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_BLUE
    p.space_after = Pt(15)
    
    bullets = [
        "• Quản lý danh mục: Thêm mới, chỉnh sửa thông tin sân bay, tàu bay, hãng bay, chuyến bay.",
        "• Điều phối chuyến bay: Cập nhật giờ bay, đổi trạng thái chuyến bay thành delayed hoặc cancelled.",
        "• Quản trị an toàn tài chính (Price Snapshot): Bảo lưu chính xác giá vé của khách hàng tại thời điểm đặt, tránh tranh chấp khi hãng hàng không điều chỉnh biểu giá mới.",
        "• Giám sát & Bảo mật: Theo dõi toàn bộ lịch sử thanh toán, xem nhật ký hoạt động (admin_logs) để phát hiện hành vi can thiệp bất thường."
    ]
    for bullet in bullets:
        p2 = tf.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_after = Pt(10)
    # ==========================================
    # SLIDE 14: Hiện trạng & Kế hoạch phát triển
    # ==========================================
    slide14 = add_content_slide("Hiện trạng & Kế hoạch phát triển")
    box_left = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.0), Inches(5.4), Inches(4.5))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = COLOR_WHITE
    box_left.line.color.rgb = COLOR_BLUE
    box_left.line.width = Pt(2)
    tf1 = box_left.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "Hiện trạng hiện tại"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_BLUE
    p.space_after = Pt(10)
    
    bullets_left = [
        "✓ Đã hoàn thiện toàn bộ schema thiết kế CSDL (19 bảng).",
        "✓ Đã xây dựng đầy đủ seed data mẫu cho các sân bay Việt Nam, tàu bay và chuyến bay thử nghiệm.",
        "✓ Đã thiết lập xong Boilerplate (cấu trúc khung) cho dự án Frontend và Backend API."
    ]
    for bullet in bullets_left:
        p2 = tf1.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_after = Pt(5)
        
    box_right = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(2.0), Inches(5.4), Inches(4.5))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = COLOR_WHITE
    box_right.line.color.rgb = COLOR_ORANGE
    box_right.line.width = Pt(2)
    tf2 = box_right.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Kế hoạch phát triển tiếp theo"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_ORANGE
    p.space_after = Pt(10)
    
    bullets_right = [
        "• Giai đoạn 1: Hiện thực hóa logic xử lý RESTful API các module backend.",
        "• Giai đoạn 2: Phát triển giao diện React chi tiết với Tailwind CSS v4.",
        "• Giai đoạn 3: Tích hợp dịch vụ thanh toán sandbox và kiểm thử tải (load testing)."
    ]
    for bullet in bullets_right:
        p2 = tf2.add_paragraph()
        p2.text = bullet
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_after = Pt(5)
    # ==========================================
    # SLIDE 15: Q&A / Cám ơn
    # ==========================================
    add_title_slide(
        "CẢM ƠN THẦY CÔ VÀ CÁC BẠN ĐÃ LẮNG NGHE!",
        "Q & A",
        "Nhóm 6\nTrường Đại học Phenikaa"
    )
    
    output_filename = r"d:\web_html_php\flight-booking-app-main\flight-booking-app-main\slide_thuyet_trinh.pptx"
    prs.save(output_filename)
    print(f"Tạo slide thuyết trình thành công tại file: {output_filename}")
if __name__ == "__main__":
    create_presentation()
